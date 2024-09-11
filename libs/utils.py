from datetime import date, datetime
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import streamlit as st

today = datetime.today()
years_ago_60 = datetime(today.year-60,1,1)

def tech_grade_calc(cr_yr, edu, certi=''):
    std_yr = [0,6,9,12]
    grd_list = ['초급','중급','고급','특급']
    std_add  = {'학사':0,'석사':-3,'전문학사':3,'박사':-9,'고졸':6}
    x = std_add.get(edu,999) #학력에 따라 기준년수 조정
    if x == 999:
        return '불가'
    std_yr = [ i + x for i in std_yr]
    # 등급 찾기
    idx = [ i  if cr_yr >= x  else  (0 if cr_yr < x else x)  for i,x in enumerate(std_yr)  ]
    grd = grd_list[max(idx)]

    if edu == '고졸' and grd=='특급':
        grd='고급'
    if grd =='초급':
        if '정보처리기사' in certi and cr_yr >= 3:
            grd='중급'
        elif '정보처리산업기사' in certi and cr_yr >= 7:
            grd='중급'

    return grd


def diff_date(d1, d2, method='months'):
    if not d1 or  str(d1) == 'NaT':
        return None
    if not d2 or  str(d2) == 'NaT':
        return None

    if method =='months':
        return int((d1.year - d2.year) * 12 + d1.month - d2.month)
    elif method == 'days':
        return int((d1 -d2).days)
    elif method == 'years':
        mm = (d1.year - d2.year) * 12 + d1.month - d2.month
        return int(mm / 12)
    elif method  == 'yearandmonth':
        mm = (d1.year - d2.year) * 12 + d1.month - d2.month
        yr, mt = divmod(mm, 12)
        return f'{yr}년{mt}개월'

def make_text_format(obj, size=10, bold=False, font_name="맑은 고딕",font_rgb=RGBColor(0,0,0)):
    obj.text_frame.paragraphs[0].font.size = Pt(size)
    obj.text_frame.paragraphs[0].font.bold = bold
    obj.text_frame.paragraphs[0].font.name = font_name
    obj.text_frame.paragraphs[0].font.color.rgb = font_rgb
    obj.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    obj.vertical_anchor = MSO_ANCHOR.MIDDLE

def edit_pres(group, data_profile, data_table, input_file, params):
    
    font_name = params.get('font_name', '맑은 고딕')
    font_size_1 = params.get('font_size_1', 11)
    font_size_2 = params.get('font_size_2', 11)
    hdr_rows = params.get('header_rows', 1)
  
  # transform input data for insertion into ppt
    table_data = data_table
  
  # get presentation template
    pres = Presentation(input_file)
  # title slide
    slide0 = pres.slides[0]
  # change title
    title = [s for s in slide0.shapes if s.has_text_frame and s.text.find("Presentation title")!=-1]
    if title:
        title[0].text = 'Title' + group

  # change subtitle
    subtitle = [s for s in slide0.shapes if s.has_text_frame and s.text.find("Subtitle")!=-1]
    if subtitle:
        subtitle[0].text = "Subtitle" + group

   # replace table text
    table = [s for s in slide0.shapes if s.has_table]
    if table:
        c = len(table[0].table.columns) #column count
        r = len(table[0].table.rows) #row count
        for i in range(0,r):
            for j in range(0,c):
                 cell = table[0].table.cell(i,j)
                 if cell.text.startswith('{') and cell.text.endswith('}'):
                    key = cell.text.replace('{','').replace('}','')
                    if key in data_profile.keys():
                        val = data_profile[key]
                        if isinstance(val, date) or isinstance(val, datetime):
                            val = val.strftime('%y.%m.%d')
                        cell.text = str(val)
                        make_text_format(cell, size=font_size_1, bold=False, font_name=font_name)

       # 헤더명으로 컬럼위치 추출 : 헤더는 n개 가능함
        c = len(table[1].table.columns) #column count
        r = len(table[1].table.rows) #row count
        pos = {}
        for r_idx in range(0, hdr_rows):
            rows = table[1].table.rows[r_idx]  #header row
            for i in range(0, c):
                pos.update({rows.cells[i].text: i})
            
        # 데이터부에 넣기: 헤더row만큼 skip
        for i in range(0, len(table_data)):
            for j in range(0, len(table_data.columns)):
                c = pos.get(table_data.columns[j], 999)
                if c < 999 and i < r-hdr_rows:
                    cell = table[1].table.cell(i+hdr_rows, c)
                    cell.text = str(table_data.iloc[i, j])
                    make_text_format(cell, size=font_size_2, bold=False, font_name=font_name)

    return pres

def split_frame(input_df, rows):
    # input df로 reindex하는 이유는 원본이 필터된 경우가 있으므로..
    df_split = [input_df.reset_index(drop=True).loc[i : i + rows - 1, :] for i in range(0, len(input_df), rows)]
    return df_split

def ini_widget(init_name, init_val):
    print('initialize widget: ', init_name, init_val)
    st.session_state[init_name] = init_val

